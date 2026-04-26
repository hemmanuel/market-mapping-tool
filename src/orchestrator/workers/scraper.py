import os
import uuid
import tempfile
import httpx
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from newspaper import Article

from src.orchestrator.core.schemas import TaskFrame, ScraperInput, ScraperOutput
from src.api.events import event_manager
from src.services.storage import storage
from src.db.session import AsyncSessionLocal
from src.models.relational import PendingDocument


def _sanitize_text(value: str) -> str:
    return value.replace("\x00", "")


class ScraperWorker:
    async def execute(self, task: TaskFrame) -> ScraperOutput:
        payload = ScraperInput(**task.payload)
        current_url = payload.url
        pipeline_id = task.pipeline_id

        raw_text = ""
        storage_object = None
        status_code = 200

        # Known problematic domains that consistently block scrapers and Jina
        SKIP_DOMAINS = [
            "bloomberg.com", "wsj.com", "forbes.com", "businesswire.com", 
            "linkedin.com", "ft.com", "sec.gov", "spglobal.com", 
            "pitchbook.com", "reuters.com", "cnbc.com", "nytimes.com"
        ]
        
        if any(domain in current_url.lower() for domain in SKIP_DOMAINS):
            await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Pre-emptively skipping known problematic domain: {current_url}"})
            return ScraperOutput(raw_text="", storage_object=None, status_code=403)
        
        await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Attempting to download: {current_url}"})
        
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
            }
            
            async with httpx.AsyncClient() as client:
                is_jina_fallback = False
                try:
                    response = await client.get(current_url, headers=headers, timeout=15.0, follow_redirects=True)
                    status_code = response.status_code
                    response.raise_for_status()
                except httpx.HTTPStatusError as e:
                    status_code = e.response.status_code
                    if status_code in (401, 403):
                        await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Encountered {status_code}. Using Jina Reader API fallback..."})
                        response = await client.get(f"https://r.jina.ai/{current_url}", headers=headers, timeout=30.0, follow_redirects=True)
                        response.raise_for_status()
                        is_jina_fallback = True
                    else:
                        raise e
                
                content_type = response.headers.get("Content-Type", "").lower()
                
                if is_jina_fallback:
                    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Extracted via Jina Reader API."})
                    raw_text = response.text
                    storage_object = f"{pipeline_id}/{uuid.uuid4()}.html"
                    storage.upload_text(raw_text, storage_object, "text/html")
                elif "application/pdf" in content_type or current_url.lower().endswith(".pdf"):
                    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Detected PDF document. Extracting text..."})
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                        tmp.write(response.content)
                        tmp_path = tmp.name
                    try:
                        storage_object = f"{pipeline_id}/{uuid.uuid4()}.pdf"
                        storage.upload_file(tmp_path, storage_object, "application/pdf")
                        reader = PdfReader(tmp_path)
                        text_parts = [page.extract_text() for page in reader.pages if page.extract_text()]
                        raw_text = " ".join(text_parts)
                    finally:
                        os.unlink(tmp_path)
                        
                elif "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in content_type or current_url.lower().endswith(".docx"):
                    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Detected Word document. Extracting text..."})
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                        tmp.write(response.content)
                        tmp_path = tmp.name
                    try:
                        storage_object = f"{pipeline_id}/{uuid.uuid4()}.docx"
                        storage.upload_file(tmp_path, storage_object, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
                        doc = Document(tmp_path)
                        raw_text = " ".join([p.text for p in doc.paragraphs])
                    finally:
                        os.unlink(tmp_path)
                        
                elif "application/vnd.openxmlformats-officedocument.presentationml.presentation" in content_type or current_url.lower().endswith(".pptx"):
                    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Detected PowerPoint document. Extracting text..."})
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                        tmp.write(response.content)
                        tmp_path = tmp.name
                    try:
                        storage_object = f"{pipeline_id}/{uuid.uuid4()}.pptx"
                        storage.upload_file(tmp_path, storage_object, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
                        prs = Presentation(tmp_path)
                        text_parts = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text_parts.append(shape.text)
                        raw_text = " ".join(text_parts)
                    finally:
                        os.unlink(tmp_path)
                        
                elif "spreadsheet" in content_type or "excel" in content_type or current_url.lower().endswith(".xlsx") or current_url.lower().endswith(".csv"):
                    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Detected Spreadsheet. Extracting data..."})
                    with tempfile.NamedTemporaryFile(delete=False) as tmp:
                        tmp.write(response.content)
                        tmp_path = tmp.name
                    try:
                        if current_url.lower().endswith(".csv") or "csv" in content_type:
                            storage_object = f"{pipeline_id}/{uuid.uuid4()}.csv"
                            storage.upload_file(tmp_path, storage_object, "text/csv")
                            df = pd.read_csv(tmp_path)
                        else:
                            storage_object = f"{pipeline_id}/{uuid.uuid4()}.xlsx"
                            storage.upload_file(tmp_path, storage_object, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
                            df = pd.read_excel(tmp_path)
                        raw_text = df.to_string()
                    finally:
                        os.unlink(tmp_path)
                        
                else:
                    # Default to HTML using newspaper3k for smart extraction
                    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Detected HTML page. Parsing with newspaper3k..."})
                    storage_object = f"{pipeline_id}/{uuid.uuid4()}.html"
                    storage.upload_text(response.text, storage_object, "text/html")
                    
                    article = Article(current_url)
                    article.set_html(response.text)
                    article.parse()
                    raw_text = article.text
                    
                    if len(raw_text) < 200:
                        await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Extracted text too short (likely paywall/boilerplate). Rejecting."})
                        raw_text = ""
                    
                # Check for massive documents
                raw_text = _sanitize_text(raw_text)
                if len(raw_text) > 500000:
                    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] WARNING: Document is exceptionally large ({len(raw_text)} chars). Skipping to prevent memory overload."})
                    await event_manager.publish(pipeline_id, {"type": "large_file_pending", "data": {"url": current_url, "estimated_size": len(raw_text)}})
                    
                    async with AsyncSessionLocal() as session:
                        try:
                            pending_doc = PendingDocument(
                                site_id=pipeline_id,
                                url=current_url,
                                estimated_size=len(raw_text),
                                status="pending"
                            )
                            session.add(pending_doc)
                            await session.commit()
                        except Exception as db_err:
                            await session.rollback()
                            print(f"Failed to save PendingDocument: {db_err}")
                    
                    raw_text = "" # Return empty text so bouncer rejects it
                else:
                    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Successfully extracted {len(raw_text)} characters."})
                
        except Exception as e:
            await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Failed to download/parse: {e}"})
            print(f"Failed to scrape {current_url}: {e}")
            raw_text = "" # Return empty text so bouncer rejects it
        
        return ScraperOutput(raw_text=raw_text, storage_object=storage_object, status_code=status_code)
