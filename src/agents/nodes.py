from typing import List, Dict, Any
import os
import time
import tempfile
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import ChatPromptTemplate
from newspaper import Article
import httpx
from pydantic import BaseModel, Field
import uuid

from src.agents.state import AgentState, SearchQuery
from src.db.neo4j_session import driver
from src.api.events import event_manager
from src.db.session import AsyncSessionLocal
from src.models.relational import DataSource, Document as PGDocument
from src.services.storage import storage
from sqlalchemy.future import select
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from dotenv import load_dotenv

load_dotenv()

from langchain_ollama import ChatOllama

# Initialize LLM and Embeddings
llm = ChatGoogleGenerativeAI(model=os.getenv("GEMINI_MODEL", "gemini-3-flash-preview"), api_key=os.getenv("GEMINI_API_KEY"))
ollama_llm = ChatOllama(model="llama3.1", base_url=os.getenv("OLLAMA_BASE_URL", "http://localhost:11434"), temperature=0)
embeddings = GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001", google_api_key=os.getenv("GEMINI_API_KEY"))

# Models for structured output
class SearchVector(BaseModel):
    intent: str = Field(description="The intent of this search vector (e.g., 'Phase 0a: Foundation - Regulatory Frameworks')")
    queries: List[str] = Field(description="List of specific search queries for this intent")
    expected_yield: str = Field(description="What kind of data this vector is expected to yield")

class PlannerOutput(BaseModel):
    search_vectors: List[SearchVector] = Field(description="List of search vectors")

class BouncerOutput(BaseModel):
    is_relevant: bool = Field(description="Whether the text is relevant to the niche")
    reason: str = Field(description="Reason for relevance or irrelevance")

class MarketSizingOutput(BaseModel):
    micro_buckets: List[str] = Field(description="List of highly specific micro-buckets to search for companies")

class CompanyExtractionOutput(BaseModel):
    companies: List[str] = Field(description="List of discovered company names")

async def market_sizing_node(state: AgentState) -> AgentState:
    """Estimates market size and generates micro-buckets for exhaustive extraction."""
    pipeline_id = state["pipeline_id"]
    niche = state["niche"]
    
    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[MarketSizing] Estimating market scale and generating micro-buckets for: {niche}"})
    
    system_prompt = """You are a Market Sizing Expert. Your goal is to help exhaustively map a market niche by breaking it down into small, highly specific "micro-buckets".
    
    First, estimate the total number of active VC-backed startups and mature companies in the niche across all stages (Seed, Series A, Series B, Late, Public).
    Then, fracture the market into specific micro-buckets (by geography, sub-niche, or vintage) such that EACH bucket likely contains fewer than 50 companies.
    
    Example micro-buckets:
    - "Seed stage Electric Power startups in North America funded in 2023"
    - "Series A Grid Storage companies in Europe"
    - "Late-stage Virtual Power Plant companies"
    
    Return a comprehensive list of these micro-buckets."""
    
    prompt = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
        ("user", "Niche: {niche}")
    ])
    
    chain = prompt | llm.with_structured_output(MarketSizingOutput)
    result = await chain.ainvoke({"niche": niche})
    
    micro_buckets = result.micro_buckets
    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[MarketSizing] Generated {len(micro_buckets)} micro-buckets for exhaustive extraction."})
    
    return {**state, "micro_buckets": micro_buckets}

import asyncio

async def company_extraction_node(state: AgentState) -> AgentState:
    """Extracts companies for each micro-bucket concurrently."""
    pipeline_id = state["pipeline_id"]
    niche = state["niche"]
    micro_buckets = state.get("micro_buckets", [])
    
    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[CompanyExtraction] Extracting companies across {len(micro_buckets)} micro-buckets..."})
    
    async def extract_companies_for_bucket(bucket: str) -> List[str]:
        # Task A: Memory-based extraction
        system_prompt = f"""You are an expert VC/PE analyst with deep knowledge of the {niche} market.
        List ALL known companies that fit perfectly into this specific micro-bucket: "{bucket}".
        Return ONLY the company names."""
        
        prompt = ChatPromptTemplate.from_messages([
            ("system", system_prompt),
            ("user", "Bucket: {bucket}")
        ])
        
        chain = prompt | llm.with_structured_output(CompanyExtractionOutput)
        memory_companies = []
        try:
            res = await chain.ainvoke({"bucket": bucket})
            memory_companies = res.companies
        except Exception as e:
            print(f"Error extracting companies for bucket {bucket}: {e}")

        # Task B: Active Search for Early-Stage
        serper_api_key = os.getenv("SERPER_API_KEY")
        search_queries = [
            f'"{bucket}" "pre-seed" OR "seed" startup',
            f'site:crunchbase.com/organization "{bucket}"'
        ]
        
        early_stage_companies = []
        if serper_api_key:
            async with httpx.AsyncClient() as client:
                for query in search_queries:
                    try:
                        response = await client.post(
                            'https://google.serper.dev/search', 
                            headers={'X-API-KEY': serper_api_key}, 
                            json={"q": query, "num": 10},
                            timeout=10.0
                        )
                        response.raise_for_status()
                        data = response.json()
                        snippets = [r.get("snippet", "") + " " + r.get("title", "") for r in data.get("organic", [])]
                        
                        if snippets:
                            extraction_prompt = ChatPromptTemplate.from_messages([
                                ("system", "Extract a list of startup company names from the following search snippets. Return ONLY the company names."),
                                ("user", "Snippets: {snippets}")
                            ])
                            chain_extract = extraction_prompt | llm.with_structured_output(CompanyExtractionOutput)
                            res_extract = await chain_extract.ainvoke({"snippets": "\n".join(snippets)})
                            early_stage_companies.extend(res_extract.companies)
                    except Exception as e:
                        print(f"Error in active search for bucket {bucket}, query {query}: {e}")

        return list(set(memory_companies + early_stage_companies))

    # Run extractions concurrently
    results = await asyncio.gather(*(extract_companies_for_bucket(bucket) for bucket in micro_buckets))
    
    # Flatten and deduplicate
    discovered_companies = list(set([company for sublist in results for company in sublist if company]))
    
    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[CompanyExtraction] Discovered {len(discovered_companies)} unique companies. Starting deep enrichment..."})
    
    from src.agents.enrichment_agent import enrich_company
    from src.agents.neo4j_enrichment import save_enriched_company_to_neo4j
    
    # Limit concurrency to avoid rate limits (Serper/Gemini)
    sem = asyncio.Semaphore(5)
    
    async def enrich_and_save(company_name: str):
        async with sem:
            try:
                enriched = await enrich_company(company_name, niche)
                if enriched:
                    await save_enriched_company_to_neo4j(pipeline_id, enriched)
            except Exception as e:
                print(f"Error enriching {company_name}: {e}")
                
    # Run enrichment concurrently
    if discovered_companies:
        await asyncio.gather(*(enrich_and_save(company) for company in discovered_companies))
        
    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[CompanyExtraction] Successfully enriched and saved companies to Neo4j."})
    
    return {**state, "discovered_companies": discovered_companies}

async def planner_node(state: AgentState) -> AgentState:
    """Generates search queries based on the niche, schema, and discovered companies."""
    pipeline_id = state["pipeline_id"]
    attempts = state.get("search_attempts", 0)
    target_urls = state.get("target_urls", 200)
    current_url_count = len(state.get("urls_to_scrape", []))
    search_feedback = state.get("search_feedback", [])
    discovered_companies = state.get("discovered_companies", [])
    micro_buckets = state.get("micro_buckets", [])
    
    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Planner] Analyzing niche: {state['niche']} (Attempt {attempts + 1})"})
    
    feedback_text = ""
    if search_feedback:
        feedback_text = "\n\nFEEDBACK FROM PREVIOUS ATTEMPTS:\n" + "\n".join(search_feedback) + "\n\nAdjust your strategy based on this feedback. If previous searches yielded low-quality data, pivot to different filetypes or domains."

    system_prompt = f"""You are a Master Sourcing Strategist. Your goal is to find the highest-density, most authoritative data sources for the user's specific market niche and ontology.

You MUST generate a structured search strategy incorporating the following elements:
1. Phase 0a: The Foundation (Top-Down): Generate queries specifically targeting broad industry overviews, major incumbents, regulatory frameworks, and standard market maps to build the core knowledge backbone.
2. Targeted Entity Deep-Dives: Generate specific search queries for the discovered companies to find primary sources (e.g., `"[Company Name]" "{state['niche']}" (funding OR technology)`).
3. Aggregator Hunting: Generate queries based on the micro-buckets to find public lists and databases (e.g., `site:crunchbase.com/organization "[Bucket]"`).

Think critically about where this specific type of data lives:
- If the user wants financial data or market sizing, target `filetype:xls`, `filetype:csv`, or SEC EDGAR.
- If the user wants supply chain or technical specs, target `filetype:pdf` (whitepapers, patents, technical manuals).
- If the user wants startup funding or M&A, target `filetype:pdf` or `filetype:pptx` (pitch decks, investor relations).
- If the user wants real-time sentiment or product launches, target standard HTML (news, press releases).

Generate highly targeted search vectors. Use advanced Google operators (`filetype:`, `site:`, `intitle:`) strategically based on the ontology.{feedback_text}"""

    prompt = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
        ("user", "Niche: {niche}\nEntities: {schema_entities}\nRelationships: {schema_relationships}\nDiscovered Companies: {discovered_companies}\nMicro-Buckets: {micro_buckets}")
    ])
    
    chain = prompt | llm.with_structured_output(PlannerOutput)
    result = await chain.ainvoke({
        "niche": state["niche"],
        "schema_entities": state["schema_entities"],
        "schema_relationships": state["schema_relationships"],
        "discovered_companies": ", ".join(discovered_companies[:50]) + ("..." if len(discovered_companies) > 50 else ""),
        "micro_buckets": ", ".join(micro_buckets)
    })
    
    flat_queries = []
    for vector in result.search_vectors:
        await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Planner] Intent: {vector.intent} (Expected: {vector.expected_yield})"})
        for q in vector.queries:
            await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Planner] Generated query: '{q}'"})
            flat_queries.append({"query": q, "target_domains": []})
    
    return {**state, "search_queries": flat_queries, "search_attempts": attempts + 1}

async def search_node(state: AgentState) -> AgentState:
    """Executes searches to find URLs using Serper.dev."""
    pipeline_id = state["pipeline_id"]
    urls = list(state.get("urls_to_scrape", []))
    
    serper_api_key = os.getenv("SERPER_API_KEY")
    if not serper_api_key:
        await event_manager.publish(pipeline_id, {"type": "log", "message": "[Searcher] Error: SERPER_API_KEY is not set."})
        return {**state, "urls_to_scrape": urls}
    
    for query in state.get("search_queries", []):
        search_term = query["query"]
        if query.get("target_domains"):
            domain_str = " OR ".join([f"site:{d}" for d in query["target_domains"]])
            search_term = f"{search_term} ({domain_str})"
            
        await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Searcher] Executing search for: '{search_term}'"})
        
        try:
            headers = {
                'X-API-KEY': serper_api_key,
                'Content-Type': 'application/json'
            }
            payload = {
                "q": search_term,
                "num": 10
            }
            
            async with httpx.AsyncClient() as client:
                response = await client.post('https://google.serper.dev/search', headers=headers, json=payload, timeout=10.0)
                response.raise_for_status()
                
                data = response.json()
                organic_results = data.get("organic", [])
                
                found_urls = 0
                for r in organic_results:
                    if "link" in r:
                        urls.append(r["link"])
                        found_urls += 1
                        
                await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Searcher] Found {found_urls} URLs for query."})
                
        except Exception as e:
            await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Searcher] Search failed for '{search_term}': {e}"})
            print(f"Serper search failed for '{search_term}': {e}")
            
    # Deduplicate URLs while preserving order
    unique_urls = []
    for url in urls:
        if url not in unique_urls:
            unique_urls.append(url)
            
    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Searcher] Total unique URLs queued: {len(unique_urls)}"})
    return {**state, "urls_to_scrape": unique_urls}

async def global_dedup_node(state: AgentState) -> AgentState:
    """Checks the database for existing URLs to avoid re-scraping and re-embedding."""
    pipeline_id = state["pipeline_id"]
    urls = state.get("urls_to_scrape", [])
    if not urls:
        return {**state, "cached_urls": []}
    
    await event_manager.publish(pipeline_id, {"type": "log", "message": "[GlobalDedup] Checking global cache for existing URLs..."})
    
    cached_urls = []
    urls_to_scrape = list(urls)
    stored_chunks = 0
    
    async with AsyncSessionLocal() as session:
        try:
            # Find which URLs exist in the DB
            result = await session.execute(
                select(PGDocument)
                .where(PGDocument.metadata_json["source_url"].astext.in_(urls))
                .order_by(PGDocument.metadata_json["source_url"].astext, PGDocument.metadata_json["chunk_index"].astext)
            )
            existing_docs = result.scalars().all()
            
            if existing_docs:
                # Group by URL
                docs_by_url = {}
                for doc in existing_docs:
                    url = doc.metadata_json.get("source_url")
                    if url:
                        if url not in docs_by_url:
                            docs_by_url[url] = []
                        docs_by_url[url].append(doc)
                
                # Process cached URLs
                if docs_by_url:
                    # Get or create data source
                    ds_result = await session.execute(select(DataSource).where(DataSource.site_id == pipeline_id).limit(1))
                    data_source = ds_result.scalars().first()
                    
                    if not data_source:
                        data_source = DataSource(
                            site_id=pipeline_id,
                            source_type="web_search",
                            name="Autonomous Web Search",
                            config={}
                        )
                        session.add(data_source)
                        await session.flush()
                    
                    for url, docs in docs_by_url.items():
                        await event_manager.publish(pipeline_id, {"type": "log", "message": f"[GlobalDedup] Cache hit for {url}. Copying {len(docs)} chunks..."})
                        for doc in docs:
                            new_doc = PGDocument(
                                data_source_id=data_source.id,
                                title=doc.title,
                                raw_text=doc.raw_text,
                                embedding=doc.embedding,
                                metadata_json=doc.metadata_json
                            )
                            session.add(new_doc)
                            stored_chunks += 1
                        
                        cached_urls.append(url)
                        if url in urls_to_scrape:
                            urls_to_scrape.remove(url)
                            
                    await session.commit()
                    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[GlobalDedup] Cloned {stored_chunks} cached chunks."})
        except Exception as e:
            await session.rollback()
            await event_manager.publish(pipeline_id, {"type": "log", "message": f"[GlobalDedup] Database error: {e}"})
            print(f"Database error in GlobalDedup: {e}")
            
    return {
        **state,
        "urls_to_scrape": urls_to_scrape,
        "cached_urls": cached_urls,
        "stored_chunks": state.get("stored_chunks", 0) + stored_chunks
    }

async def scrape_node(state: AgentState) -> AgentState:
    """Downloads text from a URL, handling multiple formats."""
    pipeline_id = state["pipeline_id"]
    urls = state.get("urls_to_scrape", [])
    if not urls:
        return {**state, "current_url": None, "raw_text": None}
    
    # Pop the first URL to process
    current_url = urls.pop(0)
    raw_text = ""
    storage_object = None
    
    # Known problematic domains that consistently block scrapers and Jina
    SKIP_DOMAINS = [
        "bloomberg.com", "wsj.com", "forbes.com", "businesswire.com", 
        "linkedin.com", "ft.com", "sec.gov", "spglobal.com", 
        "pitchbook.com", "reuters.com", "cnbc.com", "nytimes.com"
    ]
    
    if any(domain in current_url.lower() for domain in SKIP_DOMAINS):
        await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Pre-emptively skipping known problematic domain: {current_url}"})
        return {**state, "urls_to_scrape": urls, "current_url": current_url, "raw_text": "", "storage_object": None}
    
    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Attempting to download: {current_url}"})
    
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        }
        
        async with httpx.AsyncClient() as client:
            is_jina_fallback = False
            try:
                response = await client.get(current_url, headers=headers, timeout=15.0, follow_redirects=True)
                response.raise_for_status()
            except httpx.HTTPStatusError as e:
                if e.response.status_code in (401, 403):
                    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Encountered {e.response.status_code}. Using Jina Reader API fallback..."})
                    response = await client.get(f"https://r.jina.ai/{current_url}", headers=headers, timeout=30.0, follow_redirects=True)
                    response.raise_for_status()
                    is_jina_fallback = True
                else:
                    raise e
            
            content_type = response.headers.get("Content-Type", "").lower()
            
            if is_jina_fallback:
                await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Extracted via Jina Reader API."})
                raw_text = response.text
                storage_object = f"{pipeline_id}/{uuid.uuid4()}.html"
                storage.upload_text(raw_text, storage_object, "text/html")
            elif "application/pdf" in content_type or current_url.lower().endswith(".pdf"):
                await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Detected PDF document. Extracting text..."})
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                    tmp.write(response.content)
                    tmp_path = tmp.name
                try:
                    storage_object = f"{pipeline_id}/{uuid.uuid4()}.pdf"
                    storage.upload_file(tmp_path, storage_object, "application/pdf")
                    reader = PdfReader(tmp_path)
                    text_parts = [page.extract_text() for page in reader.pages if page.extract_text()]
                    raw_text = " ".join(text_parts)
                finally:
                    os.unlink(tmp_path)
                    
            elif "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in content_type or current_url.lower().endswith(".docx"):
                await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Detected Word document. Extracting text..."})
                with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                    tmp.write(response.content)
                    tmp_path = tmp.name
                try:
                    storage_object = f"{pipeline_id}/{uuid.uuid4()}.docx"
                    storage.upload_file(tmp_path, storage_object, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
                    doc = Document(tmp_path)
                    raw_text = " ".join([p.text for p in doc.paragraphs])
                finally:
                    os.unlink(tmp_path)
                    
            elif "application/vnd.openxmlformats-officedocument.presentationml.presentation" in content_type or current_url.lower().endswith(".pptx"):
                await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Detected PowerPoint document. Extracting text..."})
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                    tmp.write(response.content)
                    tmp_path = tmp.name
                try:
                    storage_object = f"{pipeline_id}/{uuid.uuid4()}.pptx"
                    storage.upload_file(tmp_path, storage_object, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
                    prs = Presentation(tmp_path)
                    text_parts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_parts.append(shape.text)
                    raw_text = " ".join(text_parts)
                finally:
                    os.unlink(tmp_path)
                    
            elif "spreadsheet" in content_type or "excel" in content_type or current_url.lower().endswith(".xlsx") or current_url.lower().endswith(".csv"):
                await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Detected Spreadsheet. Extracting data..."})
                with tempfile.NamedTemporaryFile(delete=False) as tmp:
                    tmp.write(response.content)
                    tmp_path = tmp.name
                try:
                    if current_url.lower().endswith(".csv") or "csv" in content_type:
                        storage_object = f"{pipeline_id}/{uuid.uuid4()}.csv"
                        storage.upload_file(tmp_path, storage_object, "text/csv")
                        df = pd.read_csv(tmp_path)
                    else:
                        storage_object = f"{pipeline_id}/{uuid.uuid4()}.xlsx"
                        storage.upload_file(tmp_path, storage_object, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
                        df = pd.read_excel(tmp_path)
                    raw_text = df.to_string()
                finally:
                    os.unlink(tmp_path)
                    
            else:
                # Default to HTML using newspaper3k for smart extraction
                await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Detected HTML page. Parsing with newspaper3k..."})
                storage_object = f"{pipeline_id}/{uuid.uuid4()}.html"
                storage.upload_text(response.text, storage_object, "text/html")
                
                article = Article(current_url)
                article.set_html(response.text)
                article.parse()
                raw_text = article.text
                
                if len(raw_text) < 200:
                    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Extracted text too short (likely paywall/boilerplate). Rejecting."})
                    raw_text = ""
                
            # Continue large documents through the normal ingestion path rather than
            # diverting them into a manual-review queue.
            if len(raw_text) > 500000:
                await event_manager.publish(
                    pipeline_id,
                    {
                        "type": "log",
                        "message": (
                            f"[Scraper] WARNING: Document is exceptionally large "
                            f"({len(raw_text)} chars). Continuing automatic ingestion."
                        ),
                    },
                )

            await event_manager.publish(
                pipeline_id,
                {"type": "log", "message": f"[Scraper] Successfully extracted {len(raw_text)} characters."},
            )
            
    except Exception as e:
        await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Scraper] Failed to download/parse: {e}"})
        print(f"Failed to scrape {current_url}: {e}")
        raw_text = "" # Return empty text so bouncer rejects it
    
    return {**state, "urls_to_scrape": urls, "current_url": current_url, "raw_text": raw_text, "storage_object": storage_object}

async def bouncer_node(state: AgentState) -> AgentState:
    """Scores relevance of the text using a fast, deterministic keyword density check."""
    pipeline_id = state["pipeline_id"]
    await event_manager.publish(pipeline_id, {"type": "log", "message": "[Bouncer] Evaluating text relevance with keyword density check..."})
    
    raw_text = state.get("raw_text", "")
    current_url = state.get("current_url", "unknown_url")
    search_feedback = state.get("search_feedback", [])
    
    if not raw_text:
        await event_manager.publish(pipeline_id, {"type": "log", "message": "[Bouncer] Rejected: No text to evaluate"})
        return {**state, "is_relevant": False, "relevance_reason": "No text to evaluate"}
        
    if len(raw_text) < 200:
        await event_manager.publish(pipeline_id, {"type": "log", "message": "[Bouncer] Rejected: Text too short (< 200 chars)"})
        search_feedback.append(f"URL {current_url} rejected: Text too short (< 200 chars). Likely a paywall, error page, or boilerplate.")
        return {**state, "is_relevant": False, "relevance_reason": "Text too short", "search_feedback": search_feedback}
        
    raw_text_lower = raw_text.lower()
    niche = state.get("niche", "").lower()
    entities = [e.lower() for e in state.get("schema_entities", [])]
    
    # Count occurrences of the niche and schema entities
    match_count = raw_text_lower.count(niche)
    for entity in entities:
        match_count += raw_text_lower.count(entity)
        
    if match_count < 3:
        reason = f"Low density: Found only {match_count} relevant keywords."
        await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Bouncer] Rejected: {reason}"})
        search_feedback.append(f"URL {current_url} rejected: {reason}. Try different filetypes or more specific domains.")
        return {**state, "is_relevant": False, "relevance_reason": reason, "search_feedback": search_feedback}
        
    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[Bouncer] Approved: High density (Found {match_count} relevant keywords)."})
    
    # Increment relevant URLs count
    relevant_urls_count = state.get("relevant_urls_count", 0) + 1
    
    return {**state, "is_relevant": True, "relevance_reason": f"High density ({match_count} keywords)", "relevant_urls_count": relevant_urls_count}

async def vector_storage_node(state: AgentState) -> AgentState:
    """Chunks text, generates embeddings, and saves to PostgreSQL."""
    pipeline_id = state["pipeline_id"]
    await event_manager.publish(pipeline_id, {"type": "log", "message": "[VectorStorage] Preparing to save chunks to PostgreSQL..."})
    
    if not state.get("raw_text"):
        await event_manager.publish(pipeline_id, {"type": "log", "message": "[VectorStorage] Skipped: No raw text to process."})
        return state
        
    raw_text = state["raw_text"]
    current_url = state["current_url"]
    
    # Chunk the text
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    chunks = splitter.split_text(raw_text)
    await event_manager.publish(pipeline_id, {"type": "log", "message": f"[VectorStorage] Split text into {len(chunks)} chunks."})
    
    if not chunks:
        return state
        
    # Generate embeddings
    try:
        chunk_embeddings = await embeddings.aembed_documents(chunks)
    except Exception as e:
        await event_manager.publish(pipeline_id, {"type": "log", "message": f"[VectorStorage] Failed to generate embeddings: {e}"})
        print(f"Failed to embed chunks: {e}")
        return state
        
    stored_chunks = 0
    
    async with AsyncSessionLocal() as session:
        try:
            # Find a data source for this pipeline (site)
            # This assumes at least one data source exists for the site
            result = await session.execute(select(DataSource).where(DataSource.site_id == pipeline_id).limit(1))
            data_source = result.scalars().first()
            
            if not data_source:
                data_source = DataSource(
                    site_id=pipeline_id,
                    source_type="web_search",
                    name="Autonomous Web Search",
                    config={}
                )
                session.add(data_source)
                await session.flush()
                
            for i, chunk in enumerate(chunks):
                doc = PGDocument(
                    data_source_id=data_source.id,
                    title=f"Extracted from {current_url} (Chunk {i+1})",
                    raw_text=chunk,
                    embedding=chunk_embeddings[i],
                    metadata_json={"source_url": current_url, "chunk_index": i, "storage_object": state.get("storage_object")}
                )
                session.add(doc)
                stored_chunks += 1
                
            await session.commit()
            await event_manager.publish(pipeline_id, {"type": "log", "message": f"[VectorStorage] Saved {stored_chunks} embedded chunks to PostgreSQL."})
            
            # Publish new_chunk event for the first chunk to show in the UI
            if chunks:
                await event_manager.publish(pipeline_id, {
                    "type": "new_chunk", 
                    "data": {
                        "source": current_url, 
                        "text_snippet": chunks[0][:200] + "..."
                    }
                })
            
        except Exception as e:
            await session.rollback()
            await event_manager.publish(pipeline_id, {"type": "log", "message": f"[VectorStorage] Database error: {e}"})
            print(f"Database error in VectorStorage: {e}")
            
    return {
        **state,
        "stored_chunks": state.get("stored_chunks", 0) + stored_chunks
    }
